from pptx import Presentation
import pandas as pd
import math

# Rutas a tus archivos
excel_path = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\datos.xlsx'  # Modifica esto con la ruta a tu archivo Excel
pptx_path = 'C:\\Users\\Janelka Bonilla\\Desktop\\testexcel_pptx.pptx'  # Modifica esto con la ruta donde quieres guardar la presentación

# Leer el archivo de Excel
df = pd.read_excel(excel_path)

# Crear una nueva presentación
prs = Presentation()

# Elegir el diseño de "Title and Content" usando el índice 1
slide_layout = prs.slide_layouts[1]

# Iterar sobre el DataFrame en pasos de 2
for i in range(0, len(df), 2):
    # Crear una nueva diapositiva
    slide = prs.slides.add_slide(slide_layout)
    
    # Configurar el título de la diapositiva con el valor de la columna "ACTIVIDAD"
    title = slide.shapes.title
    title.text = df.iloc[i]['ACTIVIDAD']
    
    # Si hay un siguiente valor, puedes agregarlo como contenido
    if i+1 < len(df):
        content = slide.placeholders[1]
        content.text = df.iloc[i+1]['ESTACION INICIAL'] + " hasta " + df.iloc[i+1]['ESTACION FINAL']

# Guardar la presentación
prs.save(pptx_path)