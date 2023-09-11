import math
from pptx import Presentation
import pandas as pd
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# Rutas a tus archivos
excel_path = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\datos.xlsx'  # Modifica esto con la ruta a tu archivo Excel
pptx_path = 'C:\\Users\\Janelka Bonilla\\Desktop\\testexcel_pptx.pptx'  # Modifica esto con la ruta donde quieres guardar la presentación

# Leer el archivo de Excel
df = pd.read_excel(excel_path)

# Calcula el número de diapositivas necesarias
num_slides = math.ceil(len(df) / 2)

# Crear una nueva presentación
prs = Presentation()

# Elegir el diseño en blanco usando el índice 6
slide_layout = prs.slide_layouts[6]

# Crear una diapositiva en blanco para el número calculado
for _ in range(num_slides):
    
    # Si necesitas agregar contenido a las diapositivas, este es el lugar para hacerlo
    # Añadir una diapositiva con un layout de título y contenido
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # Obtener dimensiones de la diapositiva y Convertir EMU a pulgadas
    x= Emu(prs.slide_width).inches
    y= Emu(prs.slide_height).inches

    # Añadir un título a la diapositiva----------------------------
    title = slide.shapes.title
    title.text = "Acceso 04"
    title = slide.shapes.title
    title.width = Inches(x)  # Ancho de 5 pulgadas
    title.height = Inches(0.5)  # Altura de 1 pulgada

    # Configura el formato del texto
    para = title.text_frame.paragraphs[0]
    run = para.runs[0]
    run.font.name = "Times New Roman"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)  # Color blanco
    para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Cambiar el color de fondo del cuadro de texto
    fill = title.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 32, 96)  # Esto configurará el color de fondo a rojo

    #lado Izquierdo---------------------------------------------------------------------------
    # Definir las dimensiones para la imagen
    height = Inches(3.15)
    width = Inches(3.15)
    left = Inches(4.5)-width
    top = Inches(1.8)

    #titulo 1-----------------------------------------------------------
    title1 = slide.shapes.add_textbox(left , top + height  ,width, Inches(0.6))
    text_frametitle1= title1.text_frame

    # Ajustar los márgenes del cuadro de texto
    text_frametitle1.margin_left = Inches(0)
    text_frametitle1.margin_right = Inches(0)
    text_frametitle1.margin_top = Inches(0)
    text_frametitle1.margin_bottom = Inches(0)


    t1 = text_frametitle1.add_paragraph()
    t1.text = "Excavacion en suelo"
    # Cambiar la fuente y tamaño del texto

    runt1 = t1.runs[0]
    runt1.font.name = "Times New Roman"
    runt1.font.size = Pt(16)
    t1.font.color.rgb = RGBColor(0, 0, 0)
    runt1.font.bold = True  # Hacer el texto negrita
    # Configura el formato del texto
    t1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    #imagen 1--------------------------------------------------------------------------------------------------
    # Añadir una imagen a la diapositiva
    img_path1 = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\grafo.jpg'  # Cambia esto por la ruta a tu imagen
    slide.shapes.add_picture(img_path1, left, top, height=height, width=width)

    #fecha 1-----------------------------------------------------------------------------------------------------
    # Crear una nueva forma y añadir texto
    fecha1= slide.shapes.add_textbox(left + width - Inches(1.07), top + height - Inches(0.19), Inches(1.07), Inches(0.19))
    text_framefecha1 = fecha1.text_frame

    # Ajustar los márgenes del cuadro de texto

    text_framefecha1.margin_left = Inches(0.15)
    text_framefecha1.margin_right = Inches(0)
    text_framefecha1.margin_top = Inches(0)
    text_framefecha1.margin_bottom = Inches(0.31)
    text_framefecha1.vertical_anchor = MSO_ANCHOR.MIDDLE


    f1 = text_framefecha1.add_paragraph()
    f1.text = "25-ago-23"
    # Cambiar la fuente y tamaño del texto
    runf1 = f1.runs[0]
    runf1.font.name = "Calibri (Cuerpo)"
    runf1.font.size = Pt(12)
    f1.font.color.rgb = RGBColor(0, 0, 0)
    runf1.font.bold = True  # Hacer el texto negrita
    f1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Cambiar el color de fondo del cuadro de texto
    fillfecha1 = fecha1.fill
    fillfecha1.solid()
    fillfecha1.fore_color.rgb = RGBColor(238, 236, 225)  # Esto configurará el color de fondo a rojo
    # Configurar un borde negro alrededor del cuadro de texto
    linefecha1 = fecha1.line
    linefecha1.color.rgb = RGBColor(0, 0, 0)  # Color negro
    linefecha1.width = Pt(1.25)  # Grosor de 1.25 pt

    # Añadir una descripción/subtítulo
    # Crear una nueva forma y añadir texto
    estaciones1= slide.shapes.add_textbox (left , top + height-Inches(0.3),width, Inches(1))
    text_frameestaciones1 = estaciones1.text_frame

    # Ajustar los márgenes del cuadro de texto

    text_frameestaciones1.margin_left = Inches(0)
    text_frameestaciones1.margin_right = Inches(0)
    text_frameestaciones1.margin_top = Inches(0)
    text_frameestaciones1.margin_bottom = Inches(0)


    text_frameestaciones1.clear() 
    e1 = text_frameestaciones1.add_paragraph()
    e1.text = "30+000 hasta 12+200"
    # Cambiar la fuente y tamaño del texto
    rune1 = e1.runs[0]
    rune1.font.name = "Times New Roman"
    rune1.font.size = Pt(16)
    e1.font.color.rgb = RGBColor(0, 0, 0)
    rune1.font.bold = False  # Hacer el texto negrita
    e1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    #lado derecho---------------------------------------------------------------------------
    # Definir las dimensiones para la imagen
    height = Inches(3.15)
    width = Inches(3.15)
    left = Inches(5.5)
    top = Inches(1.8)

    #titulo 2-----------------------------------------------------------
    title2 = slide.shapes.add_textbox(left , top + height  ,width, Inches(0.6))
    text_frametitle2= title2.text_frame

    # Ajustar los márgenes del cuadro de texto
    text_frametitle2.margin_left = Inches(0)
    text_frametitle2.margin_right = Inches(0)
    text_frametitle2.margin_top = Inches(0)
    text_frametitle2.margin_bottom = Inches(0)

    t2 = text_frametitle2.add_paragraph()
    t2.text = "Túnel Subterráneo"  # Cambiado el título

    # Cambiar la fuente y tamaño del texto
    runt2 = t2.runs[0]
    runt2.font.name = "Times New Roman"
    runt2.font.size = Pt(16)
    t2.font.color.rgb = RGBColor(0, 0, 0)
    runt2.font.bold = True  
    t2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    #imagen 2--------------------------------------------------------------------------------------------------
    # Añadir una imagen a la diapositiva
    img_path2 = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\grafo.jpg'  # Cambia esto por la ruta a tu nueva imagen
    slide.shapes.add_picture(img_path2, left, top, height=height, width=width)

    #fecha 2-----------------------------------------------------------------------------------------------------
    # Crear una nueva forma y añadir texto
    fecha2 = slide.shapes.add_textbox(left + width - Inches(1.07), top + height - Inches(0.19), Inches(1.07), Inches(0.19))
    text_framefecha2 = fecha2.text_frame

    # Ajustar los márgenes del cuadro de texto
    text_framefecha2.margin_left = Inches(0.15)
    text_framefecha2.margin_right = Inches(0)
    text_framefecha2.margin_top = Inches(0)
    text_framefecha2.margin_bottom = Inches(0.31)
    text_framefecha2.vertical_anchor = MSO_ANCHOR.MIDDLE

    f2 = text_framefecha2.add_paragraph()
    f2.text = "28-ago-23"  # Cambiada la fecha

    # Cambiar la fuente y tamaño del texto
    runf2 = f2.runs[0]
    runf2.font.name = "Calibri (Cuerpo)"
    runf2.font.size = Pt(12)
    f2.font.color.rgb = RGBColor(0, 0, 0)
    runf2.font.bold = True  
    f2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Cambiar el color de fondo del cuadro de texto
    fillfecha2 = fecha2.fill
    fillfecha2.solid()
    fillfecha2.fore_color.rgb = RGBColor(238, 236, 225)  
    # Configurar un borde negro alrededor del cuadro de texto
    linefecha2 = fecha2.line
    linefecha2.color.rgb = RGBColor(0, 0, 0)  
    linefecha2.width = Pt(1.25)  

    # Añadir una descripción/subtítulo
    # Crear una nueva forma y añadir texto
    estaciones2 = slide.shapes.add_textbox (left , top + height-Inches(0.3),width, Inches(1))
    text_frameestaciones2 = estaciones2.text_frame

    # Ajustar los márgenes del cuadro de texto
    text_frameestaciones2.margin_left = Inches(0)
    text_frameestaciones2.margin_right = Inches(0)
    text_frameestaciones2.margin_top = Inches(0)
    text_frameestaciones2.margin_bottom = Inches(0)

    e2 = text_frameestaciones2.add_paragraph()
    e2.text = "15+050 hasta 24+100"  # Cambiada la descripción

    # Cambiar la fuente y tamaño del texto
    rune2 = e2.runs[0]
    rune2.font.name = "Times New Roman"
    rune2.font.size = Pt(16)
    e2.font.color.rgb = RGBColor(0, 0, 0)
    rune2.font.bold = False  
    e2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER



# Guardar la presentación
prs.save(pptx_path)