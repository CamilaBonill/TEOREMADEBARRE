from pptx import Presentation

prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Prueba"

prs.save('C:\\Users\\Janelka Bonilla\\Desktop\\test_pptx.pptx')
