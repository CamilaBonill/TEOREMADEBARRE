from pptx import Presentation
import pandas as pd

# Rutas a tus archivos
excel_path = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\datos.xlsx'  # Modifica esto con la ruta a tu archivo Excel
pptx_path = 'C:\\Users\\Janelka Bonilla\\Desktop\\testexcel_pptx.pptx'  # Modifica esto con la ruta donde quieres guardar la presentación

# Leer el archivo de Excel
df = pd.read_excel(excel_path)

# Crear una nueva presentación
prs = Presentation()

# Elegir el diseño en blanco usando el índice 6
slide_layout = prs.slide_layouts[6]

# Crear una diapositiva en blanco para cada línea en el archivo de Excel
for _, row in df.iterrows():
    slide = prs.slides.add_slide(slide_layout)
    # Si necesitas agregar contenido a las diapositivas, este es el lugar para hacerlo

# Guardar la presentación
prs.save(pptx_path)