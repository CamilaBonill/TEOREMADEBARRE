
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


#img1
# Crear una nueva presentación
prs = Presentation()

# Añadir una diapositiva con un layout de título y contenido
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Obtener dimensiones de la diapositiva y Convertir EMU a pulgadas
x= Emu(prs.slide_width).inches
y= Emu(prs.slide_height).inches

# Definir las dimensiones para la imagen
height = Inches(3.15)
width = Inches(3.15)
left = Inches(0.79)
top = Inches(1.20)
x=3.15
y=3.15
l= 0.79
t=1.2
# Añadir un título a la diapositiva
title = slide.shapes.title
title.text = "Excavacion"
title_shape = slide.shapes.title
title_shape.width = width   # Ancho de 5 pulgadas
title_shape.height = Inches(1)  # Altura de 1 pulgada
title_shape.left = left   # Mover a 2 pulgadas desde el borde izquierdo
title_shape.top = top + height     # Mover a 1 pulgada desde el borde superior
title.text_frame.paragraphs[0].runs[0].font.name = "Times New Roman"
title.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Cambiar 
title.text_frame.paragraphs[0].runs[0].font.bold = True
title.text_frame.text_anchor = MSO_ANCHOR.MIDDLE

# Añadir una imagen a la diapositiva
img_path = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\grafo.jpg'  # Cambia esto por la ruta a tu imagen
slide.shapes.add_picture(img_path, left, top, height=height, width=width)


# Crear una nueva forma y añadir texto
txBox = slide.shapes.add_textbox(left + width - Inches(1.5), top + height - Inches(0.5), Inches(1.07), Inches(0.19))
text_frame = txBox.text_frame

# Ajustar los márgenes del cuadro de texto


text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

'''
# Ajuste del texto dentro del cuadro
text_frame.word_wrap = True
text_frame.auto_size = True  # Esto ajustará automáticamente el tamaño del texto para que quepa dentro del cuadro
'''
p = text_frame.add_paragraph()
p.text = "03-ago-23"
# Cambiar la fuente y tamaño del texto
run = p.runs[0]
run.font.name = "Calibri (Cuerpo)"
run.font.size = Pt(12)
p.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = True  # Hacer el texto negrita

# Cambiar el color de fondo del cuadro de texto
fill = txBox.fill
fill.solid()
fill.fore_color.rgb = RGBColor(238, 236, 225)  # Esto configurará el color de fondo a rojo
# Configurar un borde negro alrededor del cuadro de texto
line = txBox.line
line.color.rgb = RGBColor(0, 0, 0)  # Color negro
line.width = Pt(1.25)  # Grosor de 1.25 pt

'''
# Añadir una descripción/subtítulo
subtitle = slide.placeholders[1]
subtitle.text = "0+000 hasta 0+200"
#subtitle.text_frame.clear()
subtitle.width = width  # Ancho de 5 pulgadas
subtitle.height = Inches(1)  # Altura de 1 pulgada
subtitle.left = left    # Mover a 2 pulgadas desde el borde izquierdo
subtitle.top = top + Inches(3.2)    # Mover a 1 pulgada desde el borde superior
subtitle.text_frame.paragraphs[0].runs[0].font.name = "Times New Roman"
subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Cambiar 
subtitle.text_frame.text_anchor = MSO_ANCHOR.MIDDLE'''

# Crear una nueva forma y añadir texto
txBox2= slide.shapes.add_textbox (left , top + height ,width, Inches(0.5))
text_frame = txBox2.text_frame

# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


p2 = text_frame.add_paragraph()
p2.text = "0+000 hasta 0+200"
# Cambiar la fuente y tamaño del texto
run = p2.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = False  # Hacer el texto negrita

#segunda img

# Definir las dimensiones para la imagen
height = Inches(3.15)
width = Inches(3.15)
left = Inches(4)
top = Inches(1.20)
x=3.15
y=3.15
l= 0.79
t=1.2

# Añadir un título a la diapositiva

'''
title1_shape.width = width   # Ancho de 5 pulgadas
title1_shape.height = Inches(1)  # Altura de 1 pulgada
title1_shape.left = left   # Mover a 2 pulgadas desde el borde izquierdo
title1_shape.top = top + height     # Mover a 1 pulgada desde el borde superior
title1.text_frame.paragraphs[0].runs[0].font.name = "Times New Roman"
title1.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Cambiar 
title1.text_frame.paragraphs[0].runs[0].font.bold = True
title1.text_frame.text_anchor = MSO_ANCHOR.MIDDLE'''

txBox3 = slide.shapes.add_textbox(left , top + height ,width, Inches(1))
text_frame = txBox3.text_frame


# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

'''
# Ajuste del texto dentro del cuadro
text_frame.word_wrap = True
text_frame.auto_size = True  # Esto ajustará automáticamente el tamaño del texto para que quepa dentro del cuadro
'''
p3 = text_frame.add_paragraph()
p3.text = "Excavacion en suelo"
# Cambiar la fuente y tamaño del texto
run = p3.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = True  # Hacer el texto negrita
'''
# Cambiar el color de fondo del cuadro de texto
fill = txBox3.fill
fill.solid()
fill.fore_color.rgb = RGBColor(238, 236, 225)  # Esto configurará el color de fondo a rojo
# Configurar un borde negro alrededor del cuadro de texto
line = txBox.line
line.color.rgb = RGBColor(0, 0, 0)  # Color negro
line.width = Pt(1.25)  # Grosor de 1.25 pt
'''
# Añadir una imagen a la diapositiva
img_path1 = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\grafo.jpg'  # Cambia esto por la ruta a tu imagen
slide.shapes.add_picture(img_path1, left, top, height=height, width=width)


# Crear una nueva forma y añadir texto
txBox= slide.shapes.add_textbox(left + width - Inches(1.5), top + height - Inches(0.5), Inches(1.07), Inches(0.19))
text_frame = txBox.text_frame

# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

'''
# Ajuste del texto dentro del cuadro
text_frame.word_wrap = True
text_frame.auto_size = True  # Esto ajustará automáticamente el tamaño del texto para que quepa dentro del cuadro
'''
p1 = text_frame.add_paragraph()
p1.text = "25-ago-23"
# Cambiar la fuente y tamaño del texto
run = p1.runs[0]
run.font.name = "Calibri (Cuerpo)"
run.font.size = Pt(12)
p1.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = True  # Hacer el texto negrita

# Cambiar el color de fondo del cuadro de texto
fill = txBox.fill
fill.solid()
fill.fore_color.rgb = RGBColor(238, 236, 225)  # Esto configurará el color de fondo a rojo
# Configurar un borde negro alrededor del cuadro de texto
line = txBox.line
line.color.rgb = RGBColor(0, 0, 0)  # Color negro
line.width = Pt(1.25)  # Grosor de 1.25 pt


# Añadir una descripción/subtítulo
'''subtitle1 = slide.placeholders[1]
subtitle1.text = "30+000 hasta 12+200"
#subtitle.text_frame.clear()
subtitle1.width = width  # Ancho de 5 pulgadas
subtitle1.height = Inches(1)  # Altura de 1 pulgada
subtitle1.left = left         # Mover a 2 pulgadas desde el borde izquierdo
subtitle1.top = top + Inches(3.2)    # Mover a 1 pulgada desde el borde superior
subtitle1.text_frame.paragraphs[0].runs[0].font.name = "Times New Roman"
subtitle1.text_frame.paragraphs[0].runs[0].font.size = Pt(16)  # Cambiar 
subtitle1.text_frame.text_anchor = MSO_ANCHOR.MIDDLE'''

# Crear una nueva forma y añadir texto
txBox2= slide.shapes.add_textbox (left , top + height ,width, Inches(1))
text_frame = txBox2.text_frame

# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0)


text_frame.clear() 
p2 = text_frame.add_paragraph()
p2.line_spacing = 1.0  # Esto configura un espaciado de línea simple.
p2.text = "30+000 hasta 12+200"
# Cambiar la fuente y tamaño del texto
run = p2.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = False  # Hacer el texto negrita
p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p2.space_before = Pt(0)


# Guardar la presentación
prs.save('presentacion_dispositivo.pptx')

print("Presentación creada exitosamente!")




