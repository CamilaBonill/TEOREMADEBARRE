import openpyxl
from pptx import Presentation
from pptx.util import Inches

def create_presentation_from_excel(excel_path, pptx_path):
    # Cargar el archivo Excel y seleccionar la hoja activa
    workbook = openpyxl.load_workbook(excel_path)
    sheet = workbook.active
    
    # Crear una nueva presentación
    prs = Presentation()

    # Para cada fila en la hoja de Excel
    for row in sheet.iter_rows(min_row=2, values_only=True):  # asumiendo que la primera fila es el encabezado
        # Crear una nueva diapositiva
        slide_layout = prs.slide_layouts[5]  # eligiendo un layout, puedes ajustarlo según tus necesidades
        slide = prs.slides.add_slide(slide_layout)

        # Agregar el contenido de la fila de Excel a la diapositiva
        # Esto es un ejemplo, puedes ajustarlo según la estructura de tu archivo Excel
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = str(row[0])  # asumiendo que el primer valor de la fila es el título
        content.text = "\n".join(map(str, row[1:]))  # agregar el resto de los valores como contenido

    # Guardar la presentación
    prs.save(pptx_path)

# Uso
excel_path = "C:\\Users\\Janelka Bonilla\\Documents\\Github\\TEOREMADEBARRE\\PRUEBA1\\datos.xlsx"
pptx_path = "ruta_donde_guardar_pptx.pptx"
create_presentation_from_excel(excel_path, pptx_path)
