
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# Crear una nueva presentación
prs = Presentation()

# Añadir una diapositiva con un layout de título y contenido
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Obtener dimensiones de la diapositiva y Convertir EMU a pulgadas
x= Emu(prs.slide_width).inches
y= Emu(prs.slide_height).inches

# Añadir un título a la diapositiva----------------------------
title = slide.shapes.title
title.text = "Acceso 04"
title = slide.shapes.title
title.width = Inches(x)  # Ancho de 5 pulgadas
title.height = Inches(0.5)  # Altura de 1 pulgada

# Configura el formato del texto
para = title.text_frame.paragraphs[0]
run = para.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = RGBColor(255, 255, 255)  # Color blanco
para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# Cambiar el color de fondo del cuadro de texto
fill = title.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 32, 96)  # Esto configurará el color de fondo a rojo

#img1---------------------------------------------------------------------------
# Definir las dimensiones para la imagen
height = Inches(3.15)
width = Inches(3.15)
left = Inches(4.5)-width
top = Inches(1.20)

#titulo 1
title1= slide.shapes.add_textbox(left , top + height ,width, Inches(0.5))
text_frametitle1 = title1.text_frame

# Ajustar los márgenes del cuadro de texto
text_frametitle1.margin_left = Inches(0)
text_frametitle1.margin_right = Inches(0)
text_frametitle1.margin_top = Inches(0)
text_frametitle1.margin_bottom = Inches(0)
text_frametitle1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
text_frametitle1.vertical_anchor = MSO_ANCHOR.MIDDLE

t1 = text_frametitle1.add_paragraph()
t1.text = "Excavacion"

# Cambiar la fuente y tamaño del texto
para = title1.text_frame.paragraphs[0]
run = t1.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
t1.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = True  # Hacer el texto negrita
run.font.color.rgb = RGBColor(255, 255, 255)  # Color blanco
para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# Añadir una imagen a la diapositiva
img_path = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\grafo.jpg'  # Cambia esto por la ruta a tu imagen
slide.shapes.add_picture(img_path, left, top, height=height, width=width)


# Crear una nueva forma y añadir texto
txBox = slide.shapes.add_textbox(left + width - Inches(1.5), top + height - Inches(0.5), Inches(1.07), Inches(0.19))
text_frame = txBox.text_frame

# Ajustar los márgenes del cuadro de texto


text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

p = text_frame.add_paragraph()
p.text = "03-ago-23"
# Cambiar la fuente y tamaño del texto
run = p.runs[0]
run.font.name = "Calibri (Cuerpo)"
run.font.size = Pt(12)
p.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = True  # Hacer el texto negrita

# Cambiar el color de fondo del cuadro de texto
fill = txBox.fill
fill.solid()
fill.fore_color.rgb = RGBColor(238, 236, 225)  # Esto configurará el color de fondo a rojo
# Configurar un borde negro alrededor del cuadro de texto
line = txBox.line
line.color.rgb = RGBColor(0, 0, 0)  # Color negro
line.width = Pt(1.25)  # Grosor de 1.25 pt


# Crear una nueva forma y añadir texto
txBox2= slide.shapes.add_textbox (left , top + height ,width, Inches(0.5))
text_frame = txBox2.text_frame

# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


p2 = text_frame.add_paragraph()
p2.text = "0+000 hasta 0+200"
# Cambiar la fuente y tamaño del texto
run = p2.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = False  # Hacer el texto negrita

#segunda img-------------------------------------------------------------------

# Definir las dimensiones para la imagen
height = Inches(3.15)
width = Inches(3.15)
left = left = Inches(5.5)
top = Inches(1.20)
x=3.15
y=3.15
l= 0.79
t=1.2

txBox3 = slide.shapes.add_textbox(left , top + height ,width, Inches(1))
text_frame = txBox3.text_frame


# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


p3 = text_frame.add_paragraph()
p3.text = "Excavacion en suelo"
# Cambiar la fuente y tamaño del texto
run = p3.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = True  # Hacer el texto negrita

# Añadir una imagen a la diapositiva
img_path1 = 'C:\\Users\\Janelka Bonilla\\Documents\\GitHub\\TEOREMADEBARRE\\PRUEBA1\\grafo.jpg'  # Cambia esto por la ruta a tu imagen
slide.shapes.add_picture(img_path1, left, top, height=height, width=width)


# Crear una nueva forma y añadir texto
txBox= slide.shapes.add_textbox(left + width - Inches(1.5), top + height - Inches(0.5), Inches(1.07), Inches(0.19))
text_frame = txBox.text_frame

# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0.15)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0.31)
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


p1 = text_frame.add_paragraph()
p1.text = "25-ago-23"
# Cambiar la fuente y tamaño del texto
run = p1.runs[0]
run.font.name = "Calibri (Cuerpo)"
run.font.size = Pt(12)
p1.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = True  # Hacer el texto negrita

# Cambiar el color de fondo del cuadro de texto
fill = txBox.fill
fill.solid()
fill.fore_color.rgb = RGBColor(238, 236, 225)  # Esto configurará el color de fondo a rojo
# Configurar un borde negro alrededor del cuadro de texto
line = txBox.line
line.color.rgb = RGBColor(0, 0, 0)  # Color negro
line.width = Pt(1.25)  # Grosor de 1.25 pt


# Añadir una descripción/subtítulo
# Crear una nueva forma y añadir texto
txBox2= slide.shapes.add_textbox (left , top + height ,width, Inches(1))
text_frame = txBox2.text_frame

# Ajustar los márgenes del cuadro de texto

text_frame.margin_left = Inches(0)
text_frame.margin_right = Inches(0)
text_frame.margin_top = Inches(0)
text_frame.margin_bottom = Inches(0)


text_frame.clear() 
p2 = text_frame.add_paragraph()
p2.line_spacing = 1.0  # Esto configura un espaciado de línea simple.
p2.text = "30+000 hasta 12+200"
# Cambiar la fuente y tamaño del texto
run = p2.runs[0]
run.font.name = "Times New Roman"
run.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0, 0, 0)
run.font.bold = False  # Hacer el texto negrita
p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p2.space_before = Pt(0)


# Guardar la presentación
prs.save('C:\\Users\\Janelka Bonilla\\Desktop\\test_pptx.pptx')

print("Presentación creada exitosamente!")




