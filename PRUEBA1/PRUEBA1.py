from pptx import Presentation
from pptx.util import Inches, Pt

# Crear una nueva presentación
prs = Presentation()

# Diapositiva 1: Imagen y Texto
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])

# Insertar la imagen en la diapositiva 1
imagen_path_1 = 'DEMOJANELKA/778d2711a2f346b19526b0a4a268538a.jpg'
left = Inches(2)  # Agrega esta línea para definir 'left'
top = Inches(2)
width = Inches(6)
height = Inches(4)
picture_1 = slide_1.shapes.add_picture(imagen_path_1, left, top, width, height)

# Agregar texto debajo de la imagen
texto_1 = "Descripción de la imagen 1"
txBox_1 = picture_1.text_frame
p_1 = txBox_1.add_paragraph()
p_1.text = texto_1
p_1.font.size = Pt(12)

# Diapositiva 2: Imagen y Texto
slide_2 = prs.slides.add_slide(prs.slide_layouts[5])

# Insertar la imagen en la diapositiva 2
imagen_path_2 = 'DEMOJANELKA/778d2711a2f346b19526b0a4a268538a.jpg'
left = Inches(2)
top = Inches(2)
width = Inches(6)
height = Inches(4)
picture_2 = slide_2.shapes.add_picture(imagen_path_2, left, top, width, height)

# Agregar texto debajo de la imagen
texto_2 = "Descripción de la imagen 2"
txBox_2 = picture_2.text_frame
p_2 = txBox_2.add_paragraph()
p_2.text = texto_2
p_2.font.size = Pt(12)

# Guardar la presentación
nombre_presentacion = 'presentacion_dispositivo.pptx'
prs.save('c://Users//Janelka Bonilla//Documents//Github//TEOREMADEBARRE//presentacion_dispositivo.pptx')

print(f"Presentación '{nombre_presentacion}' creada exitosamente con las imágenes y texto.")
